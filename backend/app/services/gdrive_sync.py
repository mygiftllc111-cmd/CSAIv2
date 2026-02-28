"""Google Drive sync service for fetching documents."""
import asyncio
import io
import json
import logging
from typing import List, Dict, Optional

logger = logging.getLogger(__name__)

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB limit per CLAUDE.md

# Supported MIME types for text extraction
EXPORTABLE_MIME_TYPES = {
    "application/vnd.google-apps.document": "text/plain",
    "application/vnd.google-apps.spreadsheet": "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
}

DOWNLOADABLE_MIME_TYPES = {
    "text/plain",
    "text/markdown",
    "text/csv",
    "application/json",
    "text/html",
}

# Microsoft Office MIME types (require binary parsing)
OFFICE_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",   # .docx
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",         # .xlsx
    "application/vnd.openxmlformats-officedocument.presentationml.presentation", # .pptx
    "application/msword",        # .doc (best-effort)
    "application/vnd.ms-excel",  # .xls (best-effort)
    "application/vnd.ms-powerpoint",  # .ppt (best-effort)
}


async def fetch_gdrive_documents(
    service_account_json: str,
    target_folder: Optional[str] = None,
) -> List[Dict]:
    """
    Fetch documents from Google Drive using a service account.

    Args:
        service_account_json: Service account JSON key string or file path.
        target_folder: Optional folder name or ID to search in.

    Returns:
        List of dicts with "title", "content", "url", "external_id" keys.
    """
    if not service_account_json:
        logger.warning("Google service account JSON is empty, skipping sync")
        return []

    try:
        service = await _build_drive_service(service_account_json)
    except Exception as e:
        logger.error(f"Failed to build Google Drive service: {e}")
        return []

    loop = asyncio.get_event_loop()
    documents = []

    try:
        # List files
        file_list = await loop.run_in_executor(
            None, lambda: _list_files(service, target_folder)
        )
        logger.info(f"Found {len(file_list)} Google Drive files")

        for file_info in file_list:
            try:
                content = await loop.run_in_executor(
                    None, lambda f=file_info: _download_file_content(service, f)
                )
                if content and content.strip():
                    documents.append({
                        "title": file_info.get("name", "Untitled"),
                        "content": content,
                        "url": file_info.get("webViewLink", ""),
                        "external_id": file_info["id"],
                    })
            except Exception as e:
                logger.error(f"Error fetching GDrive file {file_info.get('name', '?')}: {e}")
                continue

    except Exception as e:
        logger.error(f"Google Drive sync error: {e}")

    return documents


async def _build_drive_service(service_account_json: str):
    """Build Google Drive API service from service account credentials."""
    from google.oauth2 import service_account as sa
    from googleapiclient.discovery import build

    # Parse JSON (could be a JSON string or file path)
    if service_account_json.strip().startswith("{"):
        info = json.loads(service_account_json)
    else:
        with open(service_account_json, "r") as f:
            info = json.load(f)

    credentials = sa.Credentials.from_service_account_info(
        info,
        scopes=["https://www.googleapis.com/auth/drive.readonly"],
    )

    loop = asyncio.get_event_loop()
    service = await loop.run_in_executor(
        None, lambda: build("drive", "v3", credentials=credentials)
    )
    return service


def _list_files(service, target_folder: Optional[str] = None) -> List[dict]:
    """List accessible files from Google Drive, recursing into subfolders."""
    supported = list(EXPORTABLE_MIME_TYPES.keys()) + list(DOWNLOADABLE_MIME_TYPES) + list(OFFICE_MIME_TYPES)
    mime_conditions = " or ".join(f"mimeType = '{m}'" for m in supported)
    FOLDER_MIME = "application/vnd.google-apps.folder"

    # Resolve starting folder IDs
    root_folder_ids: List[str] = []
    if target_folder:
        if len(target_folder) > 20 and " " not in target_folder:
            root_folder_ids = [target_folder]
        else:
            try:
                folder_results = service.files().list(
                    q=f"name = '{target_folder}' and mimeType = '{FOLDER_MIME}' and trashed = false",
                    fields="files(id, name)",
                    pageSize=5,
                ).execute()
                root_folder_ids = [f["id"] for f in folder_results.get("files", [])]
            except Exception as e:
                logger.error(f"Failed to find root folder '{target_folder}': {e}")
                return []
            if not root_folder_ids:
                logger.warning(f"Target folder '{target_folder}' not found in Google Drive")
                return []

    def _is_size_ok(f: dict) -> bool:
        """Google Workspace native files have no meaningful size — always allow them.
        Office files are also allowed regardless of size since we extract text locally."""
        mime = f.get("mimeType", "")
        if mime in EXPORTABLE_MIME_TYPES:
            return True  # will be exported as text, actual byte size doesn't apply
        if mime in OFFICE_MIME_TYPES:
            return True  # text extracted locally via python-pptx/docx/openpyxl
        size = int(f.get("size", 0))
        if size > MAX_FILE_SIZE:
            logger.warning(f"Skipping large file: {f['name']} ({size} bytes)")
            return False
        return True

    def _fetch_files_in_folder(folder_id: str) -> List[dict]:
        """Fetch all supported files directly in a folder (non-recursive)."""
        results = []
        page_token = None
        try:
            while True:
                resp = service.files().list(
                    q=f"'{folder_id}' in parents and ({mime_conditions}) and trashed = false",
                    fields="nextPageToken, files(id, name, mimeType, size, webViewLink)",
                    pageSize=100,
                    pageToken=page_token,
                ).execute()
                for f in resp.get("files", []):
                    if _is_size_ok(f):
                        results.append(f)
                page_token = resp.get("nextPageToken")
                if not page_token:
                    break
        except Exception as e:
            logger.warning(f"Failed to list files in folder {folder_id}: {e}")
        return results

    def _fetch_subfolders(folder_id: str) -> List[str]:
        """Return IDs of all direct subfolders."""
        try:
            resp = service.files().list(
                q=f"'{folder_id}' in parents and mimeType = '{FOLDER_MIME}' and trashed = false",
                fields="files(id, name)",
                pageSize=100,
            ).execute()
            return [f["id"] for f in resp.get("files", [])]
        except Exception as e:
            logger.warning(f"Failed to list subfolders of {folder_id}: {e}")
            return []

    all_files: List[dict] = []
    visited: set = set()

    if root_folder_ids:
        # BFS over subfolders
        queue = list(root_folder_ids)
        while queue:
            fid = queue.pop(0)
            if fid in visited:
                continue
            visited.add(fid)
            all_files.extend(_fetch_files_in_folder(fid))
            queue.extend(_fetch_subfolders(fid))
    else:
        # No folder filter: fetch all accessible supported files
        page_token = None
        try:
            while True:
                results = service.files().list(
                    q=f"({mime_conditions}) and trashed = false",
                    fields="nextPageToken, files(id, name, mimeType, size, webViewLink)",
                    pageSize=100,
                    pageToken=page_token,
                ).execute()
                for f in results.get("files", []):
                    if _is_size_ok(f):
                        all_files.append(f)
                page_token = results.get("nextPageToken")
                if not page_token:
                    break
        except Exception as e:
            logger.error(f"Failed to list files: {e}")

    logger.info(f"Found {len(all_files)} files across {len(visited)} folders")
    return all_files


def _download_file_content(service, file_info: dict) -> Optional[str]:
    """Download and extract text content from a Google Drive file."""
    mime_type = file_info.get("mimeType", "")
    file_id = file_info["id"]

    # Google Workspace files: export
    if mime_type in EXPORTABLE_MIME_TYPES:
        export_mime = EXPORTABLE_MIME_TYPES[mime_type]
        try:
            content = service.files().export(
                fileId=file_id,
                mimeType=export_mime,
            ).execute()
            if isinstance(content, bytes):
                return content.decode("utf-8", errors="replace")
            return str(content)
        except Exception as e:
            logger.error(f"Export failed for {file_info['name']}: {e}")
            return None

    # Regular text files: download
    if mime_type in DOWNLOADABLE_MIME_TYPES:
        try:
            content = service.files().get_media(fileId=file_id).execute()
            if isinstance(content, bytes):
                return content.decode("utf-8", errors="replace")
            return str(content)
        except Exception as e:
            logger.error(f"Download failed for {file_info['name']}: {e}")
            return None

    # Microsoft Office files: download binary and parse
    if mime_type in OFFICE_MIME_TYPES:
        try:
            raw = service.files().get_media(fileId=file_id).execute()
            return _extract_office_text(file_info["name"], raw)
        except Exception as e:
            logger.error(f"Office download failed for {file_info['name']}: {e}")
            return None

    return None


def _extract_office_text(filename: str, raw: bytes) -> Optional[str]:
    """Extract plain text from Office binary formats."""
    import io
    name_lower = filename.lower()

    if name_lower.endswith(".docx"):
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(io.BytesIO(raw))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            logger.warning(f"docx parse failed for {filename}: {e}")
            return None

    if name_lower.endswith(".xlsx") or name_lower.endswith(".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"[シート: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_text.strip():
                        lines.append(row_text)
            return "\n".join(lines) or None
        except Exception as e:
            logger.warning(f"xlsx parse failed for {filename}: {e}")
            return None

    if name_lower.endswith(".pptx") or name_lower.endswith(".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"[スライド {i}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text)
            return "\n".join(lines) or None
        except Exception as e:
            logger.warning(f"pptx parse failed for {filename}: {e}")
            return None

    return None
